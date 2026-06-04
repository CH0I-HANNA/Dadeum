from __future__ import annotations

from pathlib import Path
from unittest.mock import patch

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

import app.pipeline.slide_renderer as renderer_module
from app.pipeline.slide_renderer import render_pptx_slides


def _make_pptx(slide_count: int, tmp_path: Path, add_text: bool = False) -> Path:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for i in range(slide_count):
        slide = prs.slides.add_slide(blank_layout)
        if add_text:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            txBox.text_frame.text = f"Slide {i + 1} content"

    file_path = tmp_path / "test.pptx"
    prs.save(str(file_path))
    return file_path


class TestRenderPptxSlides:
    def test_returns_correct_slide_count(self, tmp_path):
        """유효한 PPTX → 반환 리스트 길이 == 슬라이드 수."""
        file_path = _make_pptx(5, tmp_path)
        result = render_pptx_slides(file_path)
        assert len(result) == 5

    def test_single_slide(self, tmp_path):
        file_path = _make_pptx(1, tmp_path)
        result = render_pptx_slides(file_path)
        assert len(result) == 1

    def test_returns_pil_image_instances(self, tmp_path):
        """각 반환값이 PIL.Image.Image 인스턴스."""
        file_path = _make_pptx(3, tmp_path)
        result = render_pptx_slides(file_path)
        for img in result:
            assert isinstance(img, Image.Image)

    def test_all_images_are_224x224(self, tmp_path):
        """각 이미지 크기가 (224, 224)."""
        file_path = _make_pptx(3, tmp_path)
        result = render_pptx_slides(file_path)
        for img in result:
            assert img.size == (224, 224)

    def test_images_are_rgb_mode(self, tmp_path):
        file_path = _make_pptx(2, tmp_path)
        result = render_pptx_slides(file_path)
        for img in result:
            assert img.mode == "RGB"

    def test_custom_size_parameter(self, tmp_path):
        """size 파라미터 변경 시 해당 크기로 반환."""
        file_path = _make_pptx(2, tmp_path)
        result = render_pptx_slides(file_path, size=128)
        for img in result:
            assert img.size == (128, 128)

    def test_render_failure_returns_white_image(self, tmp_path):
        """개별 슬라이드 렌더링 실패 시 흰색 224×224 이미지로 대체."""
        file_path = _make_pptx(3, tmp_path)

        original_render = renderer_module._render_slide
        call_count = 0

        def patched_render(slide, slide_w_emu, slide_h_emu, size):
            nonlocal call_count
            call_count += 1
            if call_count == 2:
                raise RuntimeError("Simulated render failure")
            return original_render(slide, slide_w_emu, slide_h_emu, size)

        with patch.object(renderer_module, "_render_slide", side_effect=patched_render):
            result = render_pptx_slides(file_path)

        assert len(result) == 3
        failed_img = result[1]
        assert failed_img.size == (224, 224)
        assert failed_img.mode == "RGB"
        pixels = list(failed_img.getdata())
        assert all(p == (255, 255, 255) for p in pixels)

    def test_no_exception_on_render_failure(self, tmp_path):
        """예외가 밖으로 전파되지 않는다."""
        file_path = _make_pptx(2, tmp_path)

        def always_fail(slide, slide_w_emu, slide_h_emu, size):
            raise RuntimeError("Always fails")

        with patch.object(renderer_module, "_render_slide", side_effect=always_fail):
            result = render_pptx_slides(file_path)  # must not raise

        assert len(result) == 2
        for img in result:
            assert img.size == (224, 224)

    def test_empty_slide_processes_normally(self, tmp_path):
        """빈 슬라이드(텍스트/이미지 없음)도 정상 처리."""
        file_path = _make_pptx(2, tmp_path)
        result = render_pptx_slides(file_path)
        assert len(result) == 2
        for img in result:
            assert isinstance(img, Image.Image)
            assert img.size == (224, 224)

    def test_slide_with_text_processes_normally(self, tmp_path):
        """텍스트가 있는 슬라이드도 정상 처리."""
        file_path = _make_pptx(2, tmp_path, add_text=True)
        result = render_pptx_slides(file_path)
        assert len(result) == 2
        for img in result:
            assert img.size == (224, 224)
