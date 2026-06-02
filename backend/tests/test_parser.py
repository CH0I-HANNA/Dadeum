import pytest
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.pipeline.parser import (
    ImageElement,
    SlideRaw,
    TextElement,
    parse_file,
    parse_pptx,
)
from app.core.exceptions import ParseError, PipelineError


@pytest.fixture
def simple_pptx(tmp_path: Path) -> Path:
    """2슬라이드 PPTX: 슬라이드1은 Arial 24pt 빨간 텍스트, 슬라이드2는 Calibri 18pt."""
    prs = Presentation()

    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Hello World"
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.italic = False
    run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox2 = slide2.shapes.add_textbox(Inches(2), Inches(2), Inches(3), Inches(1))
    p2 = txBox2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Test Slide"
    run2.font.name = "Calibri"
    run2.font.size = Pt(18)

    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def multi_slide_pptx(tmp_path: Path) -> Path:
    """5슬라이드 PPTX (슬라이드 수 테스트용)."""
    prs = Presentation()
    for _ in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Slide text"
    path = tmp_path / "multi.pptx"
    prs.save(str(path))
    return path


# --- 슬라이드 수 ---

def test_parse_pptx_returns_correct_slide_count(simple_pptx):
    result = parse_pptx(simple_pptx)
    assert len(result) == 2


def test_parse_pptx_multi_slide_count(multi_slide_pptx):
    result = parse_pptx(multi_slide_pptx)
    assert len(result) == 5


# --- SlideRaw 구조 ---

def test_parse_pptx_returns_slide_raw_instances(simple_pptx):
    result = parse_pptx(simple_pptx)
    for slide in result:
        assert isinstance(slide, SlideRaw)


def test_parse_pptx_slide_indices_are_zero_based(simple_pptx):
    result = parse_pptx(simple_pptx)
    assert result[0].slide_index == 0
    assert result[1].slide_index == 1


def test_parse_pptx_slide_dimensions_positive(simple_pptx):
    result = parse_pptx(simple_pptx)
    for slide in result:
        assert slide.slide_width_emu > 0
        assert slide.slide_height_emu > 0


# --- TextElement 추출 ---

def test_parse_pptx_text_element_font_family(simple_pptx):
    result = parse_pptx(simple_pptx)
    elements = result[0].text_elements
    assert len(elements) > 0
    assert elements[0].font_family == "Arial"


def test_parse_pptx_text_element_font_size(simple_pptx):
    result = parse_pptx(simple_pptx)
    elements = result[0].text_elements
    assert elements[0].font_size == pytest.approx(24.0)


def test_parse_pptx_text_element_color_rgb(simple_pptx):
    result = parse_pptx(simple_pptx)
    elements = result[0].text_elements
    assert elements[0].color_rgb == (255, 0, 0)


def test_parse_pptx_text_element_bold_flag(simple_pptx):
    result = parse_pptx(simple_pptx)
    assert result[0].text_elements[0].is_bold is True


def test_parse_pptx_text_element_text_content(simple_pptx):
    result = parse_pptx(simple_pptx)
    assert result[0].text_elements[0].text == "Hello World"


def test_parse_pptx_slide2_font_family(simple_pptx):
    result = parse_pptx(simple_pptx)
    elements = result[1].text_elements
    assert len(elements) > 0
    assert elements[0].font_family == "Calibri"


def test_parse_pptx_slide2_font_size(simple_pptx):
    result = parse_pptx(simple_pptx)
    elements = result[1].text_elements
    assert elements[0].font_size == pytest.approx(18.0)


# --- 위치 정규화 ---

def test_parse_pptx_position_values_in_range(simple_pptx):
    result = parse_pptx(simple_pptx)
    for slide in result:
        for elem in slide.text_elements:
            assert 0.0 <= elem.x <= 1.0, f"x out of range: {elem.x}"
            assert 0.0 <= elem.y <= 1.0, f"y out of range: {elem.y}"
            assert 0.0 <= elem.width <= 1.0, f"width out of range: {elem.width}"
            assert 0.0 <= elem.height <= 1.0, f"height out of range: {elem.height}"


def test_parse_pptx_textbox_position_normalized(simple_pptx):
    prs = Presentation(str(simple_pptx))
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    result = parse_pptx(simple_pptx)
    elem = result[0].text_elements[0]

    # 슬라이드1 텍스트박스: left=1in, top=1in
    expected_x = Inches(1) / slide_w
    expected_y = Inches(1) / slide_h
    assert elem.x == pytest.approx(expected_x, rel=1e-3)
    assert elem.y == pytest.approx(expected_y, rel=1e-3)


# --- 배경색 ---

def test_parse_pptx_background_is_tuple(simple_pptx):
    result = parse_pptx(simple_pptx)
    for slide in result:
        bg = slide.background_color_rgb
        assert isinstance(bg, tuple)
        assert len(bg) == 3
        assert all(0 <= c <= 255 for c in bg)


# --- parse_file 라우팅 ---

def test_parse_file_dispatches_pptx(simple_pptx):
    result = parse_file(simple_pptx)
    assert len(result) == 2


def test_parse_file_raises_on_unsupported_extension(tmp_path):
    bad = tmp_path / "test.docx"
    bad.write_bytes(b"fake")
    with pytest.raises(ParseError):
        parse_file(bad)


# --- 에러 처리 ---

def test_parse_pptx_raises_parse_error_on_corrupt_file(tmp_path):
    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"not a real pptx")
    with pytest.raises(ParseError):
        parse_pptx(bad)
