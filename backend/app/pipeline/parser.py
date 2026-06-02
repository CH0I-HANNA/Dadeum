from __future__ import annotations

from dataclasses import dataclass
from itertools import islice
from pathlib import Path

import pdfplumber
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

from app.core.exceptions import ParseError, PipelineError

_MAX_SLIDES = 50
_PT_TO_EMU = 12700  # 1pt = 12700 EMU (914400 EMU / 72pt)

_ALIGNMENT_MAP: dict = {
    PP_ALIGN.LEFT: "left",
    PP_ALIGN.CENTER: "center",
    PP_ALIGN.RIGHT: "right",
    PP_ALIGN.JUSTIFY: "left",
    PP_ALIGN.DISTRIBUTE: "center",
    PP_ALIGN.THAI_DISTRIBUTE: "unknown",
    PP_ALIGN.JUSTIFY_LOW: "left",
}


@dataclass
class TextElement:
    text: str
    font_family: str
    font_size: float
    is_bold: bool
    is_italic: bool
    color_rgb: tuple[int, int, int]
    x: float
    y: float
    width: float
    height: float
    alignment: str


@dataclass
class ImageElement:
    x: float
    y: float
    width: float
    height: float


@dataclass
class SlideRaw:
    slide_index: int
    text_elements: list[TextElement]
    image_elements: list[ImageElement]
    background_color_rgb: tuple[int, int, int]
    slide_width_emu: int
    slide_height_emu: int


def _get_alignment(alignment) -> str:
    if alignment is None:
        return "left"
    return _ALIGNMENT_MAP.get(alignment, "unknown")


def _get_font_size(run, para) -> float:
    if run.font.size is not None:
        return float(run.font.size.pt)
    if para.font.size is not None:
        return float(para.font.size.pt)
    return 18.0


def _get_font_color(run) -> tuple[int, int, int]:
    try:
        if run.font.color.type is not None:
            rgb = run.font.color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except Exception:
        pass
    return (0, 0, 0)


def _get_background_color(slide) -> tuple[int, int, int]:
    try:
        fill = slide.background.fill
        rgb = fill.fore_color.rgb
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return (255, 255, 255)


def _clamp(v: float) -> float:
    return max(0.0, min(1.0, v))


def _extract_text_elements(shape, slide_w: int, slide_h: int) -> list[TextElement]:
    if not shape.has_text_frame:
        return []

    x = _clamp((shape.left or 0) / slide_w)
    y = _clamp((shape.top or 0) / slide_h)
    w = _clamp((shape.width or 0) / slide_w)
    h = _clamp((shape.height or 0) / slide_h)

    elements = []
    for para in shape.text_frame.paragraphs:
        alignment = _get_alignment(para.alignment)
        for run in para.runs:
            if not run.text:
                continue
            font_family = run.font.name or para.font.name or "Unknown"
            elements.append(
                TextElement(
                    text=run.text,
                    font_family=font_family,
                    font_size=_get_font_size(run, para),
                    is_bold=bool(run.font.bold),
                    is_italic=bool(run.font.italic),
                    color_rgb=_get_font_color(run),
                    x=x,
                    y=y,
                    width=w,
                    height=h,
                    alignment=alignment,
                )
            )
    return elements


def parse_pptx(file_path: str | Path) -> list[SlideRaw]:
    """PPTX 파일을 파싱하여 슬라이드별 SlideRaw 목록을 반환한다."""
    file_path = Path(file_path)
    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        raise ParseError(f"파일을 열 수 없습니다: {e}") from e

    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    if slide_w == 0 or slide_h == 0:
        raise ParseError("슬라이드 크기를 확인할 수 없습니다")

    if len(prs.slides) == 0:
        raise PipelineError("슬라이드가 없습니다")

    result = []
    for i, slide in enumerate(islice(prs.slides, _MAX_SLIDES)):
        text_elements: list[TextElement] = []
        image_elements: list[ImageElement] = []

        for shape in slide.shapes:
            if (
                shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                or shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
            ):
                text_elements.extend(_extract_text_elements(shape, slide_w, slide_h))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_elements.append(
                    ImageElement(
                        x=_clamp((shape.left or 0) / slide_w),
                        y=_clamp((shape.top or 0) / slide_h),
                        width=_clamp((shape.width or 0) / slide_w),
                        height=_clamp((shape.height or 0) / slide_h),
                    )
                )

        result.append(
            SlideRaw(
                slide_index=i,
                text_elements=text_elements,
                image_elements=image_elements,
                background_color_rgb=_get_background_color(slide),
                slide_width_emu=slide_w,
                slide_height_emu=slide_h,
            )
        )

    return result


def parse_pdf(file_path: str | Path) -> list[SlideRaw]:
    """PDF 파일을 파싱하여 슬라이드별 SlideRaw 목록을 반환한다.
    PDF는 텍스트 위치/폰트 정보만 추출 가능하며, 이미지 위치는 bbox로 추정한다."""
    file_path = Path(file_path)
    try:
        with pdfplumber.open(str(file_path)) as pdf:
            if len(pdf.pages) == 0:
                raise PipelineError("슬라이드가 없습니다")

            result = []
            for i, page in enumerate(pdf.pages[:_MAX_SLIDES]):
                page_w = float(page.width)
                page_h = float(page.height)

                text_elements: list[TextElement] = []
                try:
                    words = page.extract_words(
                        x_tolerance=3,
                        y_tolerance=3,
                        extra_attrs=["fontname", "size"],
                    )
                    for word in words:
                        font_name = str(word.get("fontname") or "Unknown")
                        font_size = float(word.get("size") or 18.0)
                        text_elements.append(
                            TextElement(
                                text=word["text"],
                                font_family=font_name,
                                font_size=font_size,
                                is_bold=False,
                                is_italic=False,
                                color_rgb=(0, 0, 0),
                                x=_clamp(word["x0"] / page_w),
                                y=_clamp(word["top"] / page_h),
                                width=_clamp((word["x1"] - word["x0"]) / page_w),
                                height=_clamp((word["bottom"] - word["top"]) / page_h),
                                alignment="unknown",
                            )
                        )
                except Exception:
                    pass

                image_elements: list[ImageElement] = []
                try:
                    for img in page.images:
                        image_elements.append(
                            ImageElement(
                                x=_clamp(img["x0"] / page_w),
                                y=_clamp(img["top"] / page_h),
                                width=_clamp((img["x1"] - img["x0"]) / page_w),
                                height=_clamp((img["bottom"] - img["top"]) / page_h),
                            )
                        )
                except Exception:
                    pass

                result.append(
                    SlideRaw(
                        slide_index=i,
                        text_elements=text_elements,
                        image_elements=image_elements,
                        background_color_rgb=(255, 255, 255),
                        slide_width_emu=int(page_w * _PT_TO_EMU),
                        slide_height_emu=int(page_h * _PT_TO_EMU),
                    )
                )

            return result
    except PipelineError:
        raise
    except Exception as e:
        raise ParseError(f"PDF 파싱 실패: {e}") from e


def parse_file(file_path: str | Path) -> list[SlideRaw]:
    """확장자를 보고 parse_pptx 또는 parse_pdf를 호출한다."""
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()
    if suffix == ".pptx":
        return parse_pptx(file_path)
    elif suffix == ".pdf":
        return parse_pdf(file_path)
    else:
        raise ParseError(f"지원하지 않는 파일 형식입니다: {suffix}")
