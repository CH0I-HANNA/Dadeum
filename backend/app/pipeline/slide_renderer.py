from __future__ import annotations

import io
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

_CNN_SIZE = 224


def render_pptx_slides(file_path: Path, size: int = _CNN_SIZE) -> list[Image.Image]:
    """PPTX 파일의 각 슬라이드를 PIL Image (size×size RGB)로 렌더링한다.
    렌더링 실패한 슬라이드는 흰색 이미지로 대체한다.
    """
    file_path = Path(file_path)
    prs = Presentation(str(file_path))

    slide_w_emu = int(prs.slide_width)
    slide_h_emu = int(prs.slide_height)

    results: list[Image.Image] = []
    for slide in prs.slides:
        try:
            img = _render_slide(slide, slide_w_emu, slide_h_emu, size)
        except Exception:
            img = Image.new("RGB", (size, size), (255, 255, 255))
        results.append(img)

    return results


def _render_slide(slide, slide_w_emu: int, slide_h_emu: int, size: int) -> Image.Image:
    if slide_w_emu <= 0:
        slide_w_emu = 9144000
    if slide_h_emu <= 0:
        slide_h_emu = 5143500

    thumb_w = size
    thumb_h = max(1, round(thumb_w * slide_h_emu / slide_w_emu))
    scale = thumb_w / slide_w_emu

    bg = _get_bg_color(slide)
    img = Image.new("RGB", (thumb_w, thumb_h), color=bg)
    draw = ImageDraw.Draw(img)

    for shape in slide.shapes:
        try:
            _render_shape(shape, img, draw, scale)
        except Exception:
            pass

    return img.resize((size, size), Image.LANCZOS)


def _get_bg_color(slide) -> tuple[int, int, int]:
    try:
        fill = slide.background.fill
        rgb = fill.fore_color.rgb
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return (255, 255, 255)


def _safe_shape_fill(shape) -> tuple[int, int, int] | None:
    try:
        fill = shape.fill
        rgb = fill.fore_color.rgb
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


def _render_shape(shape, img: Image.Image, draw: ImageDraw.ImageDraw, scale: float) -> None:
    left = int((shape.left or 0) * scale)
    top = int((shape.top or 0) * scale)
    width = max(1, int((shape.width or 0) * scale))
    height = max(1, int((shape.height or 0) * scale))

    fill_color = _safe_shape_fill(shape)
    if fill_color:
        draw.rectangle([left, top, left + width, top + height], fill=fill_color)

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            blob = shape.image.blob
            sub = Image.open(io.BytesIO(blob)).convert("RGBA")
            sub = sub.resize((width, height), Image.LANCZOS)
            bg_sub = Image.new("RGBA", (width, height), (255, 255, 255, 255))
            bg_sub.paste(sub, mask=sub.split()[3])
            img.paste(bg_sub.convert("RGB"), (left, top))
        except Exception:
            draw.rectangle([left, top, left + width, top + height], fill=(200, 200, 200))
            draw.line([left, top, left + width, top + height], fill=(150, 150, 150), width=1)
            draw.line([left, top + height, left + width, top], fill=(150, 150, 150), width=1)

    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for sub_shape in shape.shapes:
                _render_shape(sub_shape, img, draw, scale)
        except Exception:
            pass

    elif shape.has_text_frame:
        if not fill_color:
            draw.rectangle([left, top, left + width, top + height], fill=(220, 220, 220))
        _render_text_frame(shape.text_frame, draw, left, top, width, height, scale)


def _render_text_frame(
    text_frame,
    draw: ImageDraw.ImageDraw,
    left: int,
    top: int,
    width: int,
    height: int,
    scale: float,
) -> None:
    font = ImageFont.load_default()
    y_cursor = top + 2
    for para in text_frame.paragraphs:
        line_text = "".join(run.text for run in para.runs).strip()
        if not line_text:
            y_cursor += 3
            continue
        if y_cursor >= top + height:
            break

        color: tuple[int, int, int] = (30, 30, 30)
        fs_pt = 18.0
        for run in para.runs:
            if run.text:
                try:
                    fs_pt = run.font.size.pt if run.font.size else 18.0
                except Exception:
                    pass
                try:
                    if run.font.color.type is not None:
                        rgb = run.font.color.rgb
                        color = (rgb[0], rgb[1], rgb[2])
                except Exception:
                    pass
                break

        fs_px = max(6, int(fs_pt * 12700 * scale))

        try:
            draw.text((left + 2, y_cursor), line_text, fill=color, font=font)
        except Exception:
            pass

        y_cursor += fs_px + 2
