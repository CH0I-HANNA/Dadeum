from __future__ import annotations

import io
from pathlib import Path

from fastapi import APIRouter, HTTPException
from fastapi.responses import Response
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.core.config import UPLOAD_DIR

router = APIRouter()

_THUMBNAIL_WIDTH = 800
_cache: dict[str, bytes] = {}

_FONT_DIRS = [
    "/System/Library/Fonts",
    "/System/Library/Fonts/Supplemental",
    "/Library/Fonts",
    str(Path.home() / "Library/Fonts"),
    "/usr/share/fonts",
]

# 한글 포함 fallback 순서 (한글 폰트 우선)
_FALLBACK_PATHS = [
    "/System/Library/Fonts/AppleSDGothicNeo.ttc",
    "/System/Library/Fonts/Supplemental/AppleGothic.ttf",
    "/Library/Fonts/NanumGothic.ttf",
    "/System/Library/Fonts/Helvetica.ttc",
    "/Library/Fonts/Arial.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
]

# 폰트 패밀리명 → 파일 경로 (첫 조회 시 빌드)
_name_cache: dict[str, str] = {}
_cache_ready = False


def _build_name_cache() -> None:
    global _cache_ready
    if _cache_ready:
        return
    import os
    from fontTools.ttLib import TTFont, TTCollection

    def _index(tt: "TTFont", path: str) -> None:
        fname = os.path.basename(path).lower()
        is_regular = not any(w in fname for w in ("italic", "bold", "oblique", "light", "thin", "heavy", "black", "medium", "condensed"))
        try:
            for rec in tt["name"].names:
                if rec.nameID in (1, 4):
                    try:
                        key = rec.toUnicode().lower().replace(" ", "").replace("-", "")
                        if key and (key not in _name_cache or is_regular):
                            _name_cache[key] = path
                    except Exception:
                        pass
        except Exception:
            pass

    for font_dir in _FONT_DIRS:
        if not os.path.isdir(font_dir):
            continue
        for root, _, files in os.walk(font_dir):
            for fname in files:
                if not fname.lower().endswith((".ttf", ".otf", ".ttc")):
                    continue
                path = os.path.join(root, fname)
                try:
                    if path.lower().endswith(".ttc"):
                        for tt in TTCollection(path):
                            _index(tt, path)
                    else:
                        _index(TTFont(path, lazy=True), path)
                except Exception:
                    stem = fname.rsplit(".", 1)[0].lower().replace(" ", "").replace("-", "")
                    _name_cache[stem] = path
    _cache_ready = True


def _find_font(name: str) -> str | None:
    _build_name_cache()
    key = name.lower().replace(" ", "").replace("-", "")
    if key in _name_cache:
        return _name_cache[key]
    for k, v in _name_cache.items():
        if key in k or k in key:
            return v
    return None


def _load_font(size: int, name: str | None = None) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    size = max(6, size)
    if name:
        path = _find_font(name)
        if path:
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                pass
    for path in _FALLBACK_PATHS:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            continue
    return ImageFont.load_default()


def _render_line(
    draw: ImageDraw.ImageDraw,
    x: int, y: int,
    text: str,
    name: str | None,
    size: int,
    color: tuple[int, int, int],
) -> int:
    """한글/비한글 구간을 분리해 렌더링. 지정 폰트가 한글을 지원하지 않으면 한글 구간만 fallback 사용."""
    base = _load_font(size, name)
    # 한글 구간용 폰트: 지정 폰트가 한글 글리프를 가지면 그대로, 없으면 fallback
    korean = _load_font(size)  # fallback (AppleSDGothicNeo 우선)

    segments: list[tuple[str, bool]] = []
    if text:
        cur, cur_ko = text[0], 0xAC00 <= ord(text[0]) <= 0xD7A3
        for ch in text[1:]:
            ko = 0xAC00 <= ord(ch) <= 0xD7A3
            if ko == cur_ko:
                cur += ch
            else:
                segments.append((cur, cur_ko))
                cur, cur_ko = ch, ko
        segments.append((cur, cur_ko))

    cx, line_h = x, size
    for seg, is_ko in segments:
        font = korean if is_ko else base
        try:
            draw.text((cx, y), seg, fill=color, font=font)
            bbox = draw.textbbox((cx, y), seg, font=font)
            line_h = max(line_h, bbox[3] - bbox[1])
            cx += bbox[2] - bbox[0]
        except Exception:
            draw.text((cx, y), seg, fill=color)
            cx += size * len(seg) // 2
    return line_h


def _get_bg_color(slide) -> tuple[int, int, int]:
    try:
        fill = slide.background.fill
        rgb = fill.fore_color.rgb
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return (255, 255, 255)


def _safe_shape_fill(shape) -> tuple[int, int, int] | None:
    try:
        fill = shape.fill
        rgb = fill.fore_color.rgb
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


def _render_shape(shape, img: Image.Image, draw: ImageDraw.ImageDraw, scale: float, thumb_h: int) -> None:
    left = int((shape.left or 0) * scale)
    top = int((shape.top or 0) * scale)
    width = max(1, int((shape.width or 0) * scale))
    height = max(1, int((shape.height or 0) * scale))

    fill_color = _safe_shape_fill(shape)
    if fill_color:
        draw.rectangle([left, top, left + width, top + height], fill=fill_color)

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            blob = shape.image.blob
            sub = Image.open(io.BytesIO(blob)).convert("RGBA")
            sub = sub.resize((width, height), Image.LANCZOS)
            bg_sub = Image.new("RGBA", (width, height), (255, 255, 255, 255))
            bg_sub.paste(sub, mask=sub.split()[3])
            img.paste(bg_sub.convert("RGB"), (left, top))
        except Exception:
            draw.rectangle([left, top, left + width, top + height], fill=(200, 200, 200))
            draw.line([left, top, left + width, top + height], fill=(150, 150, 150), width=2)
            draw.line([left, top + height, left + width, top], fill=(150, 150, 150), width=2)

    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for sub_shape in shape.shapes:
                _render_shape(sub_shape, img, draw, scale, thumb_h)
        except Exception:
            pass

    elif shape.has_text_frame:
        y_cursor = top + 2
        for para in shape.text_frame.paragraphs:
            line_text = "".join(run.text for run in para.runs)
            if not line_text.strip():
                y_cursor += 4
                continue

            fs_pt = 18.0
            color: tuple[int, int, int] = (0, 0, 0)
            font_name: str | None = None
            for run in para.runs:
                if run.text:
                    try:
                        fs_pt = run.font.size.pt if run.font.size else 18.0
                    except Exception:
                        pass
                    try:
                        if run.font.color.type is not None:
                            rgb = run.font.color.rgb
                            color = (rgb[0], rgb[1], rgb[2])
                    except Exception:
                        pass
                    try:
                        font_name = run.font.name
                    except Exception:
                        pass
                    break

            fs_px = max(6, int(fs_pt * 12700 * scale))

            if y_cursor + fs_px > top + height:
                break

            line_h = _render_line(draw, left + 4, y_cursor, line_text, font_name, fs_px, color)
            y_cursor += line_h + 2


def _render_pptx_slide(file_path: Path, slide_index: int) -> bytes:
    prs = Presentation(str(file_path))
    if slide_index >= len(prs.slides):
        raise HTTPException(status_code=404, detail="슬라이드 번호가 범위를 벗어났습니다.")

    slide_w_emu = int(prs.slide_width)
    slide_h_emu = int(prs.slide_height)
    slide = prs.slides[slide_index]

    thumb_h = round(_THUMBNAIL_WIDTH * slide_h_emu / slide_w_emu)
    scale = _THUMBNAIL_WIDTH / slide_w_emu

    bg = _get_bg_color(slide)
    img = Image.new("RGB", (_THUMBNAIL_WIDTH, thumb_h), color=bg)
    draw = ImageDraw.Draw(img)

    for shape in slide.shapes:
        try:
            _render_shape(shape, img, draw, scale, thumb_h)
        except Exception:
            pass

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _render_pdf_slide(file_path: Path, slide_index: int) -> bytes:
    import fitz

    doc = fitz.open(str(file_path))
    if slide_index >= len(doc):
        raise ValueError("슬라이드 번호가 범위를 벗어났습니다.")
    page = doc[slide_index]
    scale = _THUMBNAIL_WIDTH / page.rect.width
    mat = fitz.Matrix(scale, scale)
    pix = page.get_pixmap(matrix=mat)
    buf = io.BytesIO(pix.tobytes("png"))
    return buf.getvalue()


def _find_file_path(file_id: str) -> Path | None:
    for ext in (".pptx", ".pdf"):
        path = UPLOAD_DIR / f"{file_id}{ext}"
        if path.exists():
            return path
    return None


@router.get("/thumbnail/{file_id}/{slide_num}")
async def get_thumbnail(file_id: str, slide_num: int) -> Response:
    if slide_num < 0:
        raise HTTPException(status_code=400, detail="slide_num은 0 이상이어야 합니다.")

    file_path = _find_file_path(file_id)
    if file_path is None:
        raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")

    cache_key = f"{file_id}:{slide_num}"
    if cache_key in _cache:
        return Response(content=_cache[cache_key], media_type="image/png")

    try:
        if file_path.suffix == ".pdf":
            png_bytes = _render_pdf_slide(file_path, slide_num)
        else:
            png_bytes = _render_pptx_slide(file_path, slide_num)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"썸네일 생성 실패: {e}") from e

    _cache[cache_key] = png_bytes
    return Response(content=png_bytes, media_type="image/png")
