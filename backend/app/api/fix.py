from __future__ import annotations

import io
import os
import re
import tempfile
from pathlib import Path

import fitz  # pymupdf
from fastapi import APIRouter, HTTPException, Query
from fastapi.responses import Response
from pptx import Presentation
from pptx.dml.color import RGBColor

from app.api.thumbnail import _find_file_path, _find_font, _render_pptx_slide
from app.core import task_store
from app.models.schemas import AnalysisResult, FixRequest, OutlierSlide

router = APIRouter()


# ── PPTX fix ──────────────────────────────────────────────────────────────────

def _apply_fixes_to_slide(slide, outlier: OutlierSlide) -> None:
    target_font: str | None = None
    target_color: RGBColor | None = None

    for rc in outlier.root_causes:
        if rc.feature_group == "typography" and target_font is None:
            target_font = rc.expected_value
        if rc.feature_group == "color" and target_color is None:
            m = re.search(r"RGB\((\d+),\s*(\d+),\s*(\d+)\)", rc.expected_value)
            if m:
                target_color = RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if target_font:
                    run.font.name = target_font
                if target_color:
                    run.font.color.rgb = target_color


# ── PDF fix ────────────────────────────────────────────────────────────────────

def _parse_rgb_tuple(value: str) -> tuple[float, float, float] | None:
    """'RGB(255, 128, 0)' → (1.0, 0.5, 0.0). fitz는 0~1 float 사용."""
    m = re.search(r"RGB\((\d+),\s*(\d+),\s*(\d+)\)", value)
    if m:
        return (int(m.group(1)) / 255, int(m.group(2)) / 255, int(m.group(3)) / 255)
    return None


def _int_to_rgb(packed: int) -> tuple[float, float, float]:
    """fitz span color(packed int) → (r, g, b) 0~1."""
    return ((packed >> 16) / 255, ((packed >> 8) & 0xFF) / 255, (packed & 0xFF) / 255)


def _fix_pdf_page(page: fitz.Page, outlier: OutlierSlide) -> None:
    """한 페이지의 텍스트 span을 redact 후 수정된 폰트/색상으로 재삽입한다."""
    target_font_name: str | None = None
    target_color: tuple[float, float, float] | None = None

    for rc in outlier.root_causes:
        if rc.feature_group == "typography" and target_font_name is None:
            target_font_name = rc.expected_value
        if rc.feature_group == "color" and target_color is None:
            target_color = _parse_rgb_tuple(rc.expected_value)

    # 수정할 내용이 없으면 건너뜀
    if not target_font_name and not target_color:
        return

    # 폰트 파일 경로 확인 및 fitz에 등록
    font_alias = "fix_font"
    font_path = _find_font(target_font_name) if target_font_name else None
    if font_path:
        fitz_font = fitz.Font(fontfile=font_path)
        page.insert_font(fontname=font_alias, fontbuffer=fitz_font.buffer)
    else:
        font_alias = "helv"  # 폰트 못 찾으면 Helvetica fallback

    # 원본 span 정보 수집 (redact 전에 미리 저장)
    spans: list[dict] = []
    for block in page.get_text("dict")["blocks"]:
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                if not span["text"].strip():
                    continue
                spans.append({
                    "text": span["text"],
                    "origin": span["origin"],
                    "bbox": span["bbox"],
                    "size": span["size"],
                    "color": _int_to_rgb(span["color"]),
                })

    # redact: 모든 텍스트 영역을 배경색으로 지움
    for span in spans:
        page.add_redact_annot(fitz.Rect(span["bbox"]), fill=(1, 1, 1))
    page.apply_redactions()

    # 재삽입: 수정 대상(폰트/색상)만 교체, 나머지는 원본 유지
    for span in spans:
        color = target_color if target_color else span["color"]
        page.insert_text(
            fitz.Point(span["origin"]),
            span["text"],
            fontname=font_alias,
            fontsize=span["size"],
            color=color,
        )


def _fix_pdf(file_path: Path, result: AnalysisResult) -> bytes:
    doc = fitz.open(str(file_path))
    for outlier in result.outlier_slides:
        if outlier.slide_index < len(doc):
            _fix_pdf_page(doc[outlier.slide_index], outlier)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _render_pdf_page_png(pdf_bytes: bytes, page_index: int) -> bytes:
    doc = fitz.open(stream=pdf_bytes)
    pix = doc[page_index].get_pixmap(dpi=150)
    return pix.tobytes("png")


# ── 엔드포인트 ─────────────────────────────────────────────────────────────────

@router.post("/fix/{file_id}")
async def fix_presentation(file_id: str, task_body: FixRequest) -> Response:
    task = task_store.get_task(task_body.task_id)
    if task is None or task["status"] != "completed":
        raise HTTPException(status_code=404, detail="분석 결과를 찾을 수 없습니다.")

    result: AnalysisResult = task["result"]
    file_path = _find_file_path(file_id)
    if file_path is None:
        raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")

    if file_path.suffix == ".pptx":
        prs = Presentation(str(file_path))
        for outlier in result.outlier_slides:
            _apply_fixes_to_slide(prs.slides[outlier.slide_index], outlier)
        buf = io.BytesIO()
        prs.save(buf)
        return Response(
            content=buf.getvalue(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="dadeum-fixed-{file_id[:8]}.pptx"'},
        )

    if file_path.suffix == ".pdf":
        pdf_bytes = _fix_pdf(file_path, result)
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="dadeum-fixed-{file_id[:8]}.pdf"'},
        )

    raise HTTPException(status_code=400, detail="지원하지 않는 파일 형식입니다.")


@router.get("/preview-fix/{file_id}/{slide_num}")
async def preview_fix(file_id: str, slide_num: int, task_id: str = Query(...)) -> Response:
    if slide_num < 0:
        raise HTTPException(status_code=400, detail="slide_num은 0 이상이어야 합니다.")

    task = task_store.get_task(task_id)
    if task is None or task["status"] != "completed":
        raise HTTPException(status_code=404, detail="분석 결과를 찾을 수 없습니다.")

    result: AnalysisResult = task["result"]
    file_path = _find_file_path(file_id)
    if file_path is None:
        raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")

    outlier = next((o for o in result.outlier_slides if o.slide_index == slide_num), None)
    if outlier is None:
        raise HTTPException(status_code=404, detail="해당 슬라이드는 이상 슬라이드가 아닙니다.")

    if file_path.suffix == ".pptx":
        prs = Presentation(str(file_path))
        _apply_fixes_to_slide(prs.slides[slide_num], outlier)
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = Path(tmp.name)
            prs.save(str(tmp_path))
        try:
            png_bytes = _render_pptx_slide(tmp_path, slide_num)
        finally:
            os.unlink(tmp_path)
        return Response(content=png_bytes, media_type="image/png")

    if file_path.suffix == ".pdf":
        pdf_bytes = _fix_pdf(file_path, result)
        png_bytes = _render_pdf_page_png(pdf_bytes, slide_num)
        return Response(content=png_bytes, media_type="image/png")

    raise HTTPException(status_code=400, detail="지원하지 않는 파일 형식입니다.")
