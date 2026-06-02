"""
테스트용 PPTX 생성 스크립트.
일관성 있는 슬라이드 8장 + 이상 슬라이드 2장(색상/폰트 이탈)을 만든다.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import random

W = Inches(13.33)
H = Inches(7.5)


def add_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def set_bg(slide, r, g, b):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_text(slide, text, left, top, width, height, font_name, font_size, bold=False, color=(255, 255, 255)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return txBox


def add_rect(slide, left, top, width, height, r, g, b):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    return shape


prs = Presentation()
prs.slide_width = W
prs.slide_height = H

BG = (15, 15, 20)          # 거의 검정 배경 (일관)
ACCENT = (255, 180, 30)    # 황금색 강조 (일관)
FONT_MAIN = "Calibri"
FONT_SIZE_TITLE = 36
FONT_SIZE_BODY = 18

# ── 슬라이드 1: 표지 (정상) ──────────────────────────────
s = add_slide(prs)
set_bg(s, *BG)
add_rect(s, Inches(0), Inches(3.2), W, Pt(4), *ACCENT)
add_text(s, "다듬 디자인 가이드", Inches(1), Inches(1.5), Inches(11), Inches(1.5),
         FONT_MAIN, FONT_SIZE_TITLE, bold=True)
add_text(s, "슬라이드 일관성 분석 테스트 덱 · 2026", Inches(1), Inches(3.0), Inches(10), Inches(0.8),
         FONT_MAIN, 20, color=(180, 180, 180))

# ── 슬라이드 2: 섹션 헤더 (정상) ──────────────────────────
s = add_slide(prs)
set_bg(s, *BG)
add_rect(s, Inches(1), Inches(2.8), Inches(0.12), Inches(1.8), *ACCENT)
add_text(s, "01  개요", Inches(1.3), Inches(2.9), Inches(10), Inches(1.5),
         FONT_MAIN, FONT_SIZE_TITLE, bold=True)

# ── 슬라이드 3: 본문 (정상) ──────────────────────────────
s = add_slide(prs)
set_bg(s, *BG)
add_text(s, "핵심 목표", Inches(1), Inches(0.8), Inches(11), Inches(1),
         FONT_MAIN, FONT_SIZE_TITLE, bold=True)
body = (
    "• 슬라이드 전반의 폰트 일관성 유지\n"
    "• 색상 팔레트를 3가지 이내로 제한\n"
    "• 텍스트 영역 비율 25~45% 준수\n"
    "• 이미지와 텍스트의 균형 배치"
)
add_text(s, body, Inches(1), Inches(2.0), Inches(11), Inches(4),
         FONT_MAIN, FONT_SIZE_BODY, color=(210, 210, 210))

# ── 슬라이드 4: 본문 (정상) ──────────────────────────────
s = add_slide(prs)
set_bg(s, *BG)
add_text(s, "분석 방법론", Inches(1), Inches(0.8), Inches(11), Inches(1),
         FONT_MAIN, FONT_SIZE_TITLE, bold=True)
add_rect(s, Inches(1), Inches(2.2), Inches(3.2), Inches(2.8), 30, 30, 40)
add_rect(s, Inches(4.8), Inches(2.2), Inches(3.2), Inches(2.8), 30, 30, 40)
add_rect(s, Inches(8.6), Inches(2.2), Inches(3.2), Inches(2.8), 30, 30, 40)
add_text(s, "Feature\nExtraction", Inches(1.2), Inches(3.0), Inches(2.8), Inches(1.2),
         FONT_MAIN, FONT_SIZE_BODY, color=(210, 210, 210))
add_text(s, "Isolation\nForest", Inches(5.0), Inches(3.0), Inches(2.8), Inches(1.2),
         FONT_MAIN, FONT_SIZE_BODY, color=(210, 210, 210))
add_text(s, "Root Cause\nExplainer", Inches(8.8), Inches(3.0), Inches(2.8), Inches(1.2),
         FONT_MAIN, FONT_SIZE_BODY, color=(210, 210, 210))

# ── 슬라이드 5: 본문 (정상) ──────────────────────────────
s = add_slide(prs)
set_bg(s, *BG)
add_text(s, "성능 지표", Inches(1), Inches(0.8), Inches(11), Inches(1),
         FONT_MAIN, FONT_SIZE_TITLE, bold=True)
add_text(s,
         "정밀도(Precision)  ··········  87%\n"
         "재현율(Recall)     ··········  82%\n"
         "F1 Score           ··········  84%\n"
         "처리 속도          ··········  < 2초 / 덱",
         Inches(1.5), Inches(2.0), Inches(9), Inches(4),
         FONT_MAIN, FONT_SIZE_BODY, color=(210, 210, 210))

# ── 슬라이드 6: 본문 (정상) ──────────────────────────────
s = add_slide(prs)
set_bg(s, *BG)
add_text(s, "활용 사례", Inches(1), Inches(0.8), Inches(11), Inches(1),
         FONT_MAIN, FONT_SIZE_TITLE, bold=True)
add_text(s,
         "기업 IR 자료 · 학술 발표 · 교육 콘텐츠\n\n"
         "디자이너가 없는 스타트업이나 개인 창작자가\n"
         "프레젠테이션의 시각적 일관성을 자동으로 검사하고\n"
         "구체적인 개선 방향을 제시받을 수 있습니다.",
         Inches(1), Inches(2.0), Inches(11), Inches(4),
         FONT_MAIN, FONT_SIZE_BODY, color=(210, 210, 210))

# ── 슬라이드 7: 본문 (정상) ──────────────────────────────
s = add_slide(prs)
set_bg(s, *BG)
add_text(s, "로드맵", Inches(1), Inches(0.8), Inches(11), Inches(1),
         FONT_MAIN, FONT_SIZE_TITLE, bold=True)
add_text(s,
         "Q2 2026  ·  MVP 출시 — Isolation Forest 기반\n"
         "Q3 2026  ·  AutoEncoder 고도화\n"
         "Q4 2026  ·  GNN 기반 슬라이드 관계 분석\n"
         "Q1 2027  ·  실시간 피드백 API 제공",
         Inches(1), Inches(2.0), Inches(11), Inches(4),
         FONT_MAIN, FONT_SIZE_BODY, color=(210, 210, 210))

# ── 슬라이드 8: 본문 (정상) ──────────────────────────────
s = add_slide(prs)
set_bg(s, *BG)
add_text(s, "팀 소개", Inches(1), Inches(0.8), Inches(11), Inches(1),
         FONT_MAIN, FONT_SIZE_TITLE, bold=True)
add_text(s,
         "ML Engineer  ·  Feature 설계 및 모델 학습\n"
         "Backend Dev  ·  FastAPI 파이프라인\n"
         "Frontend Dev  ·  React 대시보드\n"
         "Designer  ·  UX 리서치 및 프로토타입",
         Inches(1), Inches(2.0), Inches(11), Inches(4),
         FONT_MAIN, FONT_SIZE_BODY, color=(210, 210, 210))

# ── 슬라이드 9: 이상 슬라이드 — 색상 이탈 ────────────────
# 배경을 완전히 다른 색상으로 변경 (흰색 계열)
s = add_slide(prs)
set_bg(s, 245, 245, 250)   # ← 흰색 배경 (BG와 완전 다름)
add_text(s, "별첨 A: 추가 데이터", Inches(1), Inches(0.8), Inches(11), Inches(1),
         FONT_MAIN, FONT_SIZE_TITLE, bold=True, color=(20, 20, 30))  # 어두운 텍스트
add_text(s,
         "이 슬라이드는 색상 팔레트에서 벗어난 이상 슬라이드입니다.\n"
         "배경색이 나머지 슬라이드와 크게 다릅니다.",
         Inches(1), Inches(2.0), Inches(11), Inches(3),
         FONT_MAIN, FONT_SIZE_BODY, color=(60, 60, 80))

# ── 슬라이드 10: 이상 슬라이드 — 폰트 이탈 ──────────────
# 폰트를 완전히 다른 것으로 변경, 크기도 이상하게
s = add_slide(prs)
set_bg(s, *BG)
add_text(s, "별첨 B: 참고 문헌", Inches(1), Inches(0.8), Inches(11), Inches(1.2),
         "Times New Roman", 48, bold=False)  # ← 다른 폰트, 다른 크기
add_text(s,
         "[1] Breunig et al. LOF: Identifying Density-Based Local Outliers. SIGMOD 2000.\n"
         "[2] Liu et al. Isolation Forest. ICDM 2008.\n"
         "[3] Kingma & Welling. Auto-Encoding Variational Bayes. ICLR 2014.",
         Inches(1), Inches(2.2), Inches(11), Inches(4),
         "Georgia", 14, color=(200, 200, 200))  # ← 또 다른 폰트

out = "/home/user/Dadeum/test_deck.pptx"
prs.save(out)
print(f"저장 완료: {out}")
print(f"총 슬라이드: {len(prs.slides)}장")
print("슬라이드 9, 10이 이상 슬라이드로 탐지될 것으로 예상")
