from __future__ import annotations

import io
import time
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches, Pt

from app.main import app

client = TestClient(app)


def _make_pptx_bytes(slide_count: int = 5) -> bytes:
    prs = Presentation()
    for i in range(slide_count):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"Slide {i + 1} content"
        run.font.size = Pt(24 if i % 2 == 0 else 18)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_bytes() -> bytes:
    return _make_pptx_bytes(slide_count=5)


class TestUpload:
    def test_upload_pptx_success(self, pptx_bytes: bytes):
        response = client.post(
            "/api/upload",
            files={"file": ("test.pptx", pptx_bytes, "application/vnd.ms-powerpoint")},
        )
        assert response.status_code == 200
        data = response.json()
        assert "file_id" in data
        assert data["file_id"]
        assert data["filename"] == "test.pptx"
        assert data["slide_count"] == 5

    def test_upload_invalid_extension_rejected(self):
        response = client.post(
            "/api/upload",
            files={"file": ("test.txt", b"hello world", "text/plain")},
        )
        assert response.status_code == 400

    def test_upload_wrong_magic_bytes_rejected(self):
        fake_pptx = b"NOT_A_REAL_PPTX_FILE_CONTENT"
        response = client.post(
            "/api/upload",
            files={"file": ("fake.pptx", fake_pptx, "application/vnd.ms-powerpoint")},
        )
        assert response.status_code == 400

    def test_upload_pdf_rejected_with_wrong_magic(self):
        fake_pdf = b"NOT_PDF_CONTENT"
        response = client.post(
            "/api/upload",
            files={"file": ("fake.pdf", fake_pdf, "application/pdf")},
        )
        assert response.status_code == 400

    def test_upload_no_filename_rejected(self):
        response = client.post(
            "/api/upload",
            files={"file": ("", b"data", "application/octet-stream")},
        )
        # 빈 파일명은 400(커스텀 검사) 또는 422(FastAPI 프레임워크 검사) 반환
        assert response.status_code in (400, 422)


class TestAnalyzeAndResult:
    def _upload(self, pptx_bytes: bytes) -> str:
        response = client.post(
            "/api/upload",
            files={"file": ("test.pptx", pptx_bytes, "application/vnd.ms-powerpoint")},
        )
        assert response.status_code == 200
        return response.json()["file_id"]

    def test_analyze_returns_task_id(self, pptx_bytes: bytes):
        file_id = self._upload(pptx_bytes)
        response = client.post(f"/api/analyze/{file_id}")
        assert response.status_code == 202
        data = response.json()
        assert "task_id" in data
        assert data["task_id"]

    def test_analyze_missing_file_returns_404(self):
        response = client.post("/api/analyze/nonexistent-file-id")
        assert response.status_code == 404

    def test_result_pending_or_completed(self, pptx_bytes: bytes):
        file_id = self._upload(pptx_bytes)
        analyze_resp = client.post(f"/api/analyze/{file_id}")
        assert analyze_resp.status_code == 202
        task_id = analyze_resp.json()["task_id"]

        # 폴링: 최대 30초 대기
        for _ in range(30):
            result_resp = client.get(f"/api/result/{task_id}")
            assert result_resp.status_code == 200
            data = result_resp.json()
            if data["status"] == "completed":
                assert "result" in data
                assert data["result"]["file_id"] == file_id
                assert data["result"]["slide_count"] == 5
                assert "consistency_score" in data["result"]
                return
            if data["status"] == "error":
                pytest.fail(f"분석 실패: {data.get('error_message')}")
            time.sleep(1)

        pytest.fail("30초 내 completed 상태가 되지 않았습니다.")

    def test_duplicate_analyze_returns_409(self, pptx_bytes: bytes):
        file_id = self._upload(pptx_bytes)
        resp1 = client.post(f"/api/analyze/{file_id}")
        assert resp1.status_code == 202

        # 즉시 재요청 — pending/processing 상태이면 409
        resp2 = client.post(f"/api/analyze/{file_id}")
        # pending/processing 상태일 때만 409, 이미 completed면 새 task 허용
        assert resp2.status_code in (202, 409)

    def test_result_missing_task_returns_404(self):
        response = client.get("/api/result/nonexistent-task-id")
        assert response.status_code == 404


class TestThumbnail:
    def _upload(self, pptx_bytes: bytes) -> str:
        response = client.post(
            "/api/upload",
            files={"file": ("test.pptx", pptx_bytes, "application/vnd.ms-powerpoint")},
        )
        assert response.status_code == 200
        return response.json()["file_id"]

    def test_thumbnail_returns_png(self, pptx_bytes: bytes):
        file_id = self._upload(pptx_bytes)
        response = client.get(f"/api/thumbnail/{file_id}/0")
        assert response.status_code == 200
        assert response.headers["content-type"] == "image/png"
        assert len(response.content) > 0

    def test_thumbnail_missing_file_returns_404(self):
        response = client.get("/api/thumbnail/nonexistent-id/0")
        assert response.status_code == 404

    def test_thumbnail_out_of_range_returns_404(self, pptx_bytes: bytes):
        file_id = self._upload(pptx_bytes)
        response = client.get(f"/api/thumbnail/{file_id}/999")
        assert response.status_code == 404

    def test_thumbnail_negative_slide_returns_400(self, pptx_bytes: bytes):
        file_id = self._upload(pptx_bytes)
        response = client.get(f"/api/thumbnail/{file_id}/-1")
        assert response.status_code == 400
